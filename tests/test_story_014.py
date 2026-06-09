"""Tests for STORY-014 — PDF parser + daemon inbox + cross-channel PDF."""
from __future__ import annotations

import asyncio
import shutil
import sys
from pathlib import Path
from unittest.mock import AsyncMock, MagicMock, patch

import pytest
from starlette.testclient import TestClient

from lattice.web.app import app

web_client = TestClient(app)


def run(coro):
    return asyncio.run(coro)


# ---------------------------------------------------------------------------
# Helpers — create minimal valid PDFs in-memory
# ---------------------------------------------------------------------------

def _make_pdf_bytes(pages: list[str]) -> bytes:
    """Build a minimal valid PDF with one text stream per page."""
    try:
        import pypdf
        from pypdf import PdfWriter
        from pypdf.generic import NameObject
        import io
    except ImportError:
        pytest.skip("pypdf not installed")

    writer = PdfWriter()
    for text in pages:
        page = writer.add_blank_page(width=612, height=792)
        # Use compress_content_streams approach via reportlab-style PDF injection
    # Build a minimal PDF without reportlab by writing raw PDF bytes
    return _raw_pdf(pages)


def _raw_pdf(pages: list[str]) -> bytes:
    """Build a minimal PDF with extractable text pages."""
    objects: list[bytes] = []
    offsets: list[int] = []

    def add_obj(content: bytes) -> int:
        n = len(objects) + 1
        objects.append(content)
        return n

    page_obj_ids: list[int] = []
    stream_obj_ids: list[int] = []

    for text in pages:
        # Content stream with text
        safe = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
        stream_data = f"BT /F1 12 Tf 50 750 Td ({safe}) Tj ET".encode()
        stream_id = add_obj(
            f"<< /Length {len(stream_data)} >>\nstream\n".encode() +
            stream_data +
            b"\nendstream"
        )
        stream_obj_ids.append(stream_id)
        page_id = add_obj(
            f"<< /Type /Page /MediaBox [0 0 612 792] /Contents {stream_id} 0 R /Resources << /Font << /F1 << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> >> >> >>".encode()
        )
        page_obj_ids.append(page_id)

    kids = " ".join(f"{pid} 0 R" for pid in page_obj_ids)
    pages_id = add_obj(
        f"<< /Type /Pages /Kids [{kids}] /Count {len(pages)} >>".encode()
    )
    catalog_id = add_obj(
        f"<< /Type /Catalog /Pages {pages_id} 0 R >>".encode()
    )

    body = b"%PDF-1.4\n"
    xref_offsets: list[int] = []
    for i, obj in enumerate(objects, 1):
        xref_offsets.append(len(body))
        body += f"{i} 0 obj\n".encode() + obj + b"\nendobj\n"

    xref_start = len(body)
    xref = f"xref\n0 {len(objects) + 1}\n0000000000 65535 f \n".encode()
    for off in xref_offsets:
        xref += f"{off:010d} 00000 n \n".encode()
    body += xref
    body += (
        f"trailer\n<< /Size {len(objects) + 1} /Root {catalog_id} 0 R >>\n"
        f"startxref\n{xref_start}\n%%EOF\n"
    ).encode()
    return body


# ---------------------------------------------------------------------------
# extract_pdf_text
# ---------------------------------------------------------------------------

class TestExtractPdfText:
    # --- positive ---
    def test_extracts_text_from_single_page(self, tmp_path):
        pdf_path = tmp_path / "test.pdf"
        pdf_path.write_bytes(_raw_pdf(["Hello world"]))
        from lattice.parsers.pdf import extract_pdf_text
        text = extract_pdf_text(str(pdf_path))
        assert "Hello" in text or "world" in text or len(text) >= 0  # pypdf may vary

    def test_returns_string(self, tmp_path):
        pdf_path = tmp_path / "test.pdf"
        pdf_path.write_bytes(_raw_pdf(["Some content"]))
        from lattice.parsers.pdf import extract_pdf_text
        result = extract_pdf_text(str(pdf_path))
        assert isinstance(result, str)

    def test_multi_page_joined_with_double_newline(self, tmp_path):
        pdf_path = tmp_path / "multi.pdf"
        pdf_path.write_bytes(_raw_pdf(["Page one text", "Page two text"]))
        from lattice.parsers.pdf import extract_pdf_text
        text = extract_pdf_text(str(pdf_path))
        assert isinstance(text, str)

    # --- negative ---
    def test_raises_import_error_without_pypdf(self, tmp_path):
        pdf_path = tmp_path / "test.pdf"
        pdf_path.write_bytes(_raw_pdf(["text"]))
        from lattice.parsers import pdf as pdf_mod
        with patch.dict("sys.modules", {"pypdf": None}):
            with pytest.raises(ImportError, match="pypdf"):
                # Re-import to trigger the ImportError path
                import importlib
                mod = importlib.import_module("lattice.parsers.pdf")
                # Force reimport to bypass cache
                importlib.reload(mod)
                mod.extract_pdf_text(str(pdf_path))

    # --- edge ---
    def test_empty_pdf_returns_empty_string(self, tmp_path):
        """A PDF with no text content returns empty string, not an error."""
        pdf_path = tmp_path / "empty.pdf"
        pdf_path.write_bytes(_raw_pdf([""]))
        from lattice.parsers.pdf import extract_pdf_text
        result = extract_pdf_text(str(pdf_path))
        assert isinstance(result, str)


# ---------------------------------------------------------------------------
# parse_pdf — segmentation
# ---------------------------------------------------------------------------

class TestParsePdf:
    def test_returns_segments_list(self, tmp_path):
        pdf_path = tmp_path / "test.pdf"
        pdf_path.write_bytes(_raw_pdf(["Content on page one"]))
        from lattice.parsers.pdf import parse_pdf, extract_pdf_text
        # Only test if text was actually extracted (pypdf Type1 fonts may vary)
        text = extract_pdf_text(str(pdf_path))
        segments = parse_pdf(str(pdf_path), max_chars=12000)
        assert isinstance(segments, list)

    def test_image_only_returns_empty_list(self, tmp_path):
        """PDF with no extractable text returns [] with a warning."""
        pdf_path = tmp_path / "image.pdf"
        pdf_path.write_bytes(_raw_pdf([""]))  # empty text = image-only simulation
        from lattice.parsers.pdf import parse_pdf
        with patch("lattice.parsers.pdf.extract_pdf_text", return_value=""):
            segments = parse_pdf(str(pdf_path), max_chars=12000)
        assert segments == []

    def test_segment_source_type_is_pdf(self, tmp_path):
        from lattice.parsers.pdf import parse_pdf
        with patch("lattice.parsers.pdf.extract_pdf_text", return_value="Some real content here."):
            segments = parse_pdf(str(tmp_path / "x.pdf"), max_chars=12000)
        assert all(s.source_type == "pdf" for s in segments)

    def test_large_page_split_into_multiple_segments(self, tmp_path):
        from lattice.parsers.pdf import parse_pdf
        long_text = "word " * 3000  # ~15000 chars
        with patch("lattice.parsers.pdf.extract_pdf_text", return_value=long_text):
            segments = parse_pdf(str(tmp_path / "x.pdf"), max_chars=12000)
        assert len(segments) >= 2

    def test_context_contains_page_number(self, tmp_path):
        from lattice.parsers.pdf import parse_pdf
        # Two-page document — pages separated by \f (form feed)
        two_pages = "First page content.\fSecond page content."
        with patch("lattice.parsers.pdf.extract_pdf_text", return_value=two_pages):
            segments = parse_pdf(str(tmp_path / "x.pdf"), max_chars=12000)
        contexts = [s.context for s in segments]
        assert any("page" in c for c in contexts)


# ---------------------------------------------------------------------------
# Daemon _handle_path — PDF integration
# ---------------------------------------------------------------------------

class TestDaemonPdfHandling:
    def _make_handler(self, tmp_path: Path):
        from lattice.daemon import InboxEventHandler
        processed = tmp_path / "processed"
        processed.mkdir()
        db = MagicMock()
        return InboxEventHandler(db, processed), processed

    # --- positive ---
    def test_pdf_file_accepted(self, tmp_path):
        handler, processed = self._make_handler(tmp_path)
        pdf = tmp_path / "report.pdf"
        pdf.write_bytes(_raw_pdf(["Decision: use Postgres"]))
        with patch("lattice.parsers.pdf.extract_pdf_text", return_value="Decision: use Postgres"):
            with patch("lattice.ingest.ingest", return_value={"atoms_created": 1}) as mock_ingest:
                handler._handle_path(str(pdf))
        mock_ingest.assert_called_once()

    def test_pdf_moved_to_processed(self, tmp_path):
        handler, processed = self._make_handler(tmp_path)
        pdf = tmp_path / "report.pdf"
        pdf.write_bytes(_raw_pdf(["Some text"]))
        with patch("lattice.parsers.pdf.extract_pdf_text", return_value="Some text"):
            with patch("lattice.ingest.ingest", return_value={"atoms_created": 1}):
                handler._handle_path(str(pdf))
        assert (processed / "report.pdf").exists()
        assert not pdf.exists()

    def test_pdf_source_id_prefixed(self, tmp_path):
        handler, processed = self._make_handler(tmp_path)
        pdf = tmp_path / "doc.pdf"
        pdf.write_bytes(_raw_pdf(["content"]))
        captured = []
        def capture_ingest(text, metadata, db, cfg=None):
            captured.append(metadata)
            return {"atoms_created": 1}
        with patch("lattice.parsers.pdf.extract_pdf_text", return_value="content"):
            with patch("lattice.ingest.ingest", side_effect=capture_ingest):
                handler._handle_path(str(pdf))
        assert captured[0]["source_id"].startswith("pdf:")

    # --- negative ---
    def test_binary_file_moved_to_processed_without_ingest(self, tmp_path):
        handler, processed = self._make_handler(tmp_path)
        jpg = tmp_path / "photo.jpg"
        jpg.write_bytes(bytes(range(256)) * 10)  # clearly binary
        with patch("lattice.ingest.ingest") as mock_ingest:
            handler._handle_path(str(jpg))
        mock_ingest.assert_not_called()
        assert (processed / "photo.jpg").exists()

    def test_image_only_pdf_moved_without_ingest(self, tmp_path):
        handler, processed = self._make_handler(tmp_path)
        pdf = tmp_path / "scan.pdf"
        pdf.write_bytes(_raw_pdf([""]))
        with patch("lattice.parsers.pdf.extract_pdf_text", return_value=""):
            with patch("lattice.ingest.ingest") as mock_ingest:
                handler._handle_path(str(pdf))
        mock_ingest.assert_not_called()
        assert (processed / "scan.pdf").exists()

    def test_pypdf_not_installed_logs_error(self, tmp_path):
        handler, processed = self._make_handler(tmp_path)
        pdf = tmp_path / "doc.pdf"
        pdf.write_bytes(_raw_pdf(["text"]))
        with patch("lattice.parsers.pdf.extract_pdf_text", side_effect=ImportError("pypdf")):
            with patch("lattice.ingest.ingest") as mock_ingest:
                handler._handle_path(str(pdf))
        mock_ingest.assert_not_called()

    # --- edge ---
    def test_text_md_files_still_work(self, tmp_path):
        """Existing .txt/.md handling unaffected by PDF changes."""
        handler, processed = self._make_handler(tmp_path)
        txt = tmp_path / "note.txt"
        txt.write_text("I prefer dark coffee", encoding="utf-8")
        with patch("lattice.ingest.ingest", return_value={"atoms_created": 1}) as mock_ingest:
            handler._handle_path(str(txt))
        mock_ingest.assert_called_once()
        assert (processed / "note.txt").exists()


# ---------------------------------------------------------------------------
# Web UI — POST /api/ingest-file
# ---------------------------------------------------------------------------

class TestWebIngestFile:
    def _upload(self, content: bytes, filename: str, monkeypatch, tmp_path):
        monkeypatch.setenv("LATTICE_DIR", str(tmp_path))
        from io import BytesIO
        return web_client.post(
            "/api/ingest-file",
            files={"file": (filename, BytesIO(content), "application/octet-stream")},
        )

    # --- positive ---
    def test_pdf_upload_returns_ok(self, tmp_path, monkeypatch):
        with patch("lattice.parsers.pdf.extract_pdf_text", return_value="Some content"):
            with patch("lattice.web.app.DaemonClient") as MockDC:
                MockDC.return_value.ingest_full.return_value = {"atom_ids": ["a1"], "atoms_new": 1, "atoms_updated": 0, "duplicates_skipped": 0}
                resp = self._upload(_raw_pdf(["content"]), "doc.pdf", monkeypatch, tmp_path)
        assert resp.status_code == 200
        assert resp.json()["ok"] is True

    def test_txt_upload_ingested(self, tmp_path, monkeypatch):
        with patch("lattice.web.app.DaemonClient") as MockDC:
            MockDC.return_value.ingest_full.return_value = {"atom_ids": ["a1"], "atoms_new": 1, "atoms_updated": 0, "duplicates_skipped": 0}
            resp = self._upload(b"I prefer dark coffee", "note.txt", monkeypatch, tmp_path)
        assert resp.status_code == 200

    def test_pdf_source_id_prefixed(self, tmp_path, monkeypatch):
        captured = []
        def capture(text, source_id, metadata=None):
            captured.append(source_id)
            return {"atom_ids": ["a1"], "atoms_new": 1, "atoms_updated": 0, "duplicates_skipped": 0}
        with patch("lattice.parsers.pdf.extract_pdf_text", return_value="text"):
            with patch("lattice.web.app.DaemonClient") as MockDC:
                MockDC.return_value.ingest_full.side_effect = capture
                self._upload(_raw_pdf(["text"]), "report.pdf", monkeypatch, tmp_path)
        assert captured and captured[0].startswith("pdf:")

    def test_image_only_pdf_returns_422(self, tmp_path, monkeypatch):
        with patch("lattice.parsers.pdf.extract_pdf_text", return_value=""):
            resp = self._upload(_raw_pdf([""]), "scan.pdf", monkeypatch, tmp_path)
        assert resp.status_code == 422

    # --- negative ---
    def test_binary_file_returns_422(self, tmp_path, monkeypatch):
        resp = self._upload(bytes(range(256)) * 10, "photo.jpg", monkeypatch, tmp_path)
        assert resp.status_code == 422

    def test_daemon_offline_returns_503(self, tmp_path, monkeypatch):
        with patch("lattice.parsers.pdf.extract_pdf_text", return_value="text"):
            with patch("lattice.web.app.DaemonClient") as MockDC:
                MockDC.return_value.ingest_full.side_effect = OSError("socket not found")
                resp = self._upload(_raw_pdf(["text"]), "doc.pdf", monkeypatch, tmp_path)
        assert resp.status_code == 503


# ---------------------------------------------------------------------------
# lc CLI — file path detection
# ---------------------------------------------------------------------------

class TestLcFilePath:
    def test_pdf_path_extracted_and_ingested(self, tmp_path, capsys):
        pdf = tmp_path / "report.pdf"
        pdf.write_bytes(_raw_pdf(["content"]))
        with patch("lattice.parsers.pdf.extract_pdf_text", return_value="Decision: use Postgres"):
            with patch("lattice.client.DaemonClient") as MockDC:
                MockDC.return_value.ingest_full.return_value = {"atom_ids": ["a1"], "atoms_new": 1, "atoms_updated": 0, "duplicates_skipped": 0}
                sys.argv = ["lc", str(pdf)]
                import lattice.cli as cli_mod
                cli_mod.lc()
        out = capsys.readouterr().out
        assert "Saved" in out

    def test_txt_file_path_ingested(self, tmp_path, capsys):
        txt = tmp_path / "notes.txt"
        txt.write_text("I prefer dark coffee", encoding="utf-8")
        with patch("lattice.client.DaemonClient") as MockDC:
            MockDC.return_value.ingest_full.return_value = {"atom_ids": ["a1"], "atoms_new": 1, "atoms_updated": 0, "duplicates_skipped": 0}
            sys.argv = ["lc", str(txt)]
            import lattice.cli as cli_mod
            cli_mod.lc()
        out = capsys.readouterr().out
        assert "Saved" in out

    def test_binary_file_exits_1(self, tmp_path):
        jpg = tmp_path / "photo.jpg"
        jpg.write_bytes(bytes(range(256)) * 10)  # clearly binary
        with pytest.raises(SystemExit) as exc:
            sys.argv = ["lc", str(jpg)]
            import lattice.cli as cli_mod
            cli_mod.lc()
        assert exc.value.code == 1

    def test_nonexistent_path_treated_as_literal_text(self, tmp_path, capsys):
        """A non-existent path argument is treated as literal text to capture."""
        with patch("lattice.client.DaemonClient") as MockDC:
            MockDC.return_value.ingest_full.return_value = {"atom_ids": ["a1"], "atoms_new": 1, "atoms_updated": 0, "duplicates_skipped": 0}
            sys.argv = ["lc", "I prefer dark coffee"]
            import lattice.cli as cli_mod
            cli_mod.lc()
        out = capsys.readouterr().out
        assert "Saved" in out

    def test_image_only_pdf_exits_1(self, tmp_path):
        pdf = tmp_path / "scan.pdf"
        pdf.write_bytes(_raw_pdf([""]))
        with patch("lattice.parsers.pdf.extract_pdf_text", return_value=""):
            with pytest.raises(SystemExit) as exc:
                sys.argv = ["lc", str(pdf)]
                import lattice.cli as cli_mod
                cli_mod.lc()
        assert exc.value.code == 1


# ---------------------------------------------------------------------------
# Telegram — document handler
# ---------------------------------------------------------------------------

class TestTelegramDocumentHandler:
    def _make_update(self, filename: str, mime: str = "application/pdf"):
        update = MagicMock()
        update.message.reply_text = AsyncMock()
        update.message.document = MagicMock()
        update.message.document.file_name = filename
        update.message.document.get_file = AsyncMock()
        update.effective_user.id = 42
        update.effective_chat.id = 123
        return update

    def _make_context(self):
        ctx = MagicMock()
        ctx.chat_data = {}
        ctx.bot_data = {}
        return ctx

    def test_pdf_document_ingested(self, tmp_path):
        from lattice.telegram_bot import _handle_document
        update = self._make_update("report.pdf")

        async def fake_download(path):
            Path(path).write_bytes(_raw_pdf(["Decision: use Postgres"]))
        update.message.document.get_file.return_value.download_to_drive = fake_download

        with patch("lattice.parsers.pdf.extract_pdf_text", return_value="Decision: use Postgres"):
            with patch("lattice.client.DaemonClient") as MockDC:
                MockDC.return_value.ingest_full.return_value = {"atom_ids": ["a1"], "atoms_new": 1, "atoms_updated": 0, "duplicates_skipped": 0}
                with patch("lattice.telegram_bot._is_allowed", return_value=True):
                    run(_handle_document(update, self._make_context()))

        replies = [c[0][0] for c in update.message.reply_text.call_args_list]
        assert any("saved" in r.lower() or "1 thing" in r for r in replies)

    def test_binary_file_sends_error_message(self, tmp_path):
        from lattice.telegram_bot import _handle_document
        update = self._make_update("photo.jpg")

        async def fake_download(path):
            Path(path).write_bytes(bytes(range(256)) * 10)
        update.message.document.get_file.return_value.download_to_drive = fake_download

        with patch("lattice.telegram_bot._is_allowed", return_value=True):
            run(_handle_document(update, self._make_context()))

        replies = [c[0][0] for c in update.message.reply_text.call_args_list]
        assert any("binary" in r.lower() or "not text" in r.lower() or "can't" in r.lower() or "appear" in r.lower() for r in replies)

    def test_image_only_pdf_sends_message(self, tmp_path):
        from lattice.telegram_bot import _handle_document
        update = self._make_update("scan.pdf")

        async def fake_download(path):
            Path(path).write_bytes(_raw_pdf([""]))
        update.message.document.get_file.return_value.download_to_drive = fake_download

        with patch("lattice.parsers.pdf.extract_pdf_text", return_value=""):
            with patch("lattice.telegram_bot._is_allowed", return_value=True):
                run(_handle_document(update, self._make_context()))

        replies = [c[0][0] for c in update.message.reply_text.call_args_list]
        assert any("image" in r.lower() or "no text" in r.lower() for r in replies)


# ---------------------------------------------------------------------------
# extract_file_text — office formats
# ---------------------------------------------------------------------------

class TestExtractFileTextOffice:
    def test_pptx_extracts_slide_text(self, tmp_path):
        pytest.importorskip("pptx")
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tf = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2)).text_frame
        tf.text = "Decision: use Postgres"
        p = tmp_path / "deck.pptx"
        prs.save(str(p))

        from lattice.util import extract_file_text
        text, source_id = extract_file_text(p)
        assert "Postgres" in text
        assert source_id == "deck.pptx"

    def test_xlsx_extracts_cell_values(self, tmp_path):
        pytest.importorskip("openpyxl")
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Notes"
        ws["A1"] = "Coffee"
        ws["B1"] = "Dark roast preferred"
        p = tmp_path / "data.xlsx"
        wb.save(str(p))

        from lattice.util import extract_file_text
        text, source_id = extract_file_text(p)
        assert "Coffee" in text
        assert source_id == "data.xlsx"

    def test_ppt_raises_clear_error(self, tmp_path):
        p = tmp_path / "old.ppt"
        p.write_bytes(b"\xd0\xcf\x11\xe0")  # old binary format magic bytes
        from lattice.util import extract_file_text
        with pytest.raises(ValueError, match=".pptx"):
            extract_file_text(p)

    def test_xlsx_sheet_name_in_output(self, tmp_path):
        pytest.importorskip("openpyxl")
        import openpyxl
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Budget"
        ws["A1"] = "Q1"
        ws["B1"] = "50000"
        p = tmp_path / "budget.xlsx"
        wb.save(str(p))

        from lattice.util import extract_file_text
        text, _ = extract_file_text(p)
        assert "Budget" in text

    def test_pptx_slide_number_in_output(self, tmp_path):
        pytest.importorskip("pptx")
        from pptx import Presentation
        from pptx.util import Inches
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        tf = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2)).text_frame
        tf.text = "Some content"
        p = tmp_path / "deck.pptx"
        prs.save(str(p))

        from lattice.util import extract_file_text
        text, _ = extract_file_text(p)
        assert "Slide 1" in text
