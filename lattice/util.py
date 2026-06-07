from __future__ import annotations

import json
import os
import re
import tempfile
from pathlib import Path


def _normalized_subject(subject: str) -> str:
    return " ".join(re.findall(r"[a-z0-9]+", subject.lower()))


def write_file_atomic(path: Path, text: str) -> None:
    """Write *text* to *path* atomically via a temp-file rename."""
    path.parent.mkdir(parents=True, exist_ok=True)
    fd, tmp_name = tempfile.mkstemp(dir=path.parent, prefix=f".{path.name}.", text=True)
    try:
        with os.fdopen(fd, "w", encoding="utf-8") as f:
            f.write(text)
        Path(tmp_name).replace(path)
    finally:
        tmp = Path(tmp_name)
        if tmp.exists():
            tmp.unlink()


def _write_json_atomic(path: Path, data: dict) -> None:
    write_file_atomic(path, json.dumps(data))


def extract_file_text(path: str | Path) -> tuple[str, str]:
    """Extract plain text from any file. Returns (text, source_id).

    Dispatch:
    - .pdf  → pypdf (raises ImportError if not installed)
    - .docx → python-docx (raises ImportError if not installed)
    - anything else → UTF-8 text read

    Raises:
    - ImportError  if a required optional dep is missing
    - ValueError   if the file is binary / no text extractable
    - OSError      if the file cannot be read
    """
    p = Path(path)
    suffix = p.suffix.lower()

    if suffix == ".pdf":
        from lattice.parsers.pdf import extract_pdf_text
        text = extract_pdf_text(str(p))
        if not text.strip():
            raise ValueError("No text found — PDF may be image-only")
        return text, f"pdf:{p.name}"

    if suffix == ".docx":
        try:
            import docx as _docx
        except ImportError:
            raise ImportError(
                "python-docx is required for .docx files. Install with: uv sync --group docx"
            )
        doc = _docx.Document(str(p))
        _HEADING_MAP = {
            "Heading 1": "#", "Heading 2": "##", "Heading 3": "###",
            "Heading 4": "####", "Title": "#", "Subtitle": "##",
        }
        parts = []
        for para in doc.paragraphs:
            if not para.text.strip():
                continue
            style_name = para.style.name if para.style else ""
            prefix = _HEADING_MAP.get(style_name, "")
            parts.append(f"{prefix} {para.text}" if prefix else para.text)
        text = "\n\n".join(parts)
        if not text.strip():
            raise ValueError("No text found in .docx file")
        return text, p.name

    if suffix == ".pptx":
        try:
            from pptx import Presentation as _Presentation
        except ImportError:
            raise ImportError(
                "python-pptx is required for .pptx files. Install with: uv sync --group office"
            )
        prs = _Presentation(str(p))
        slides: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            parts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
            if parts:
                slides.append(f"[Slide {i}]\n" + "\n".join(parts))
        text = "\n\n".join(slides)
        if not text.strip():
            raise ValueError("No text found in .pptx file — slides may be image-only")
        return text, p.name

    if suffix in {".xlsx"}:
        try:
            import openpyxl as _openpyxl
        except ImportError:
            raise ImportError(
                "openpyxl is required for .xlsx files. Install with: uv sync --group office"
            )
        wb = _openpyxl.load_workbook(str(p), data_only=True)
        sheets: list[str] = []
        for sheet in wb.worksheets:
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None and str(c).strip()]
                if cells:
                    rows.append("\t".join(cells))
            if rows:
                sheets.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
        text = "\n\n".join(sheets)
        if not text.strip():
            raise ValueError("No text found in .xlsx file — may be empty")
        return text, p.name

    if suffix in {".xls"}:
        try:
            import xlrd as _xlrd
        except ImportError:
            raise ImportError(
                "xlrd is required for .xls files. Install with: uv sync --group office"
            )
        wb = _xlrd.open_workbook(str(p))
        sheets: list[str] = []
        for sheet in wb.sheets():
            rows: list[str] = []
            for r in range(sheet.nrows):
                cells = [str(sheet.cell_value(r, c)) for c in range(sheet.ncols)
                         if str(sheet.cell_value(r, c)).strip()]
                if cells:
                    rows.append("\t".join(cells))
            if rows:
                sheets.append(f"[Sheet: {sheet.name}]\n" + "\n".join(rows))
        text = "\n\n".join(sheets)
        if not text.strip():
            raise ValueError("No text found in .xls file — may be empty")
        return text, p.name

    if suffix == ".ppt":
        raise ValueError(
            ".ppt (old PowerPoint binary) is not supported. Save as .pptx and try again."
        )

    # Everything else: try UTF-8 text read
    try:
        text = p.read_text(encoding="utf-8", errors="strict")
    except UnicodeDecodeError:
        # Try with replacement chars — if result is >30% replacement chars it's binary
        text = p.read_text(encoding="utf-8", errors="replace")
        replacement_ratio = text.count("�") / max(len(text), 1)
        if replacement_ratio > 0.1:
            raise ValueError(f"File appears to be binary (not text-readable): {p.name}")
    if not text.strip():
        raise ValueError(f"File is empty: {p.name}")
    return text, p.name
